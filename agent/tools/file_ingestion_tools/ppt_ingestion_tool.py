from pptx import Presentation
import os
import base64
from typing import Dict, Any
from google.adk.tools import ToolContext

def ppt_ingestion_tool(ppt_path: str, tool_context: ToolContext) -> Dict[str, Any]:
    """
    Extracts text, tables, notes, and images from a PPTX file into a structured dictionary.
    Images are saved in image_dir and also returned as base64 if needed.

        Args:
        ppt_path (str): Path to the input ppt.
        tool_context (ToolContext): The tool context containing state information.
        
    Returns:
        Dict[str, Any]: A dictionary containing:
            - status: Success or Failure
    """
    try:
        image_dir = "/Users/p0s08o6/Desktop/projects/Startup Analyst/startup_analyst_agent/agent/output"
        prs = Presentation(ppt_path)
        os.makedirs(image_dir, exist_ok=True)

        extracted = {"slides": []}

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_data = {
                "slide_number": slide_idx,
                "text": [],
                "tables": [],
                "notes": "",
                "images": []
            }

            # Extract text & tables
            for shape in slide.shapes:
                # Text
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_data["text"].append(text)

                # Tables
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    slide_data["tables"].append(table_data)

                # Images
                if shape.shape_type == 13:  # PICTURE
                    image = shape.image
                    img_extension = image.ext
                    img_filename = f"slide{slide_idx}_{shape.shape_id}.{img_extension}"
                    img_path = os.path.join(image_dir, img_filename)
                    with open(img_path, "wb") as f:
                        f.write(image.blob)
                    # Optionally encode as base64
                    slide_data["images"].append({
                        "filename": img_filename,
                        "path": img_path,
                        "base64": base64.b64encode(image.blob).decode("utf-8")
                    })

            # Extract notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                slide_data["notes"] = slide.notes_slide.notes_text_frame.text.strip()

            extracted["slides"].append(slide_data)

            tool_context.state['ppt_data'] = extracted

        return {
            "status": "success",
        }

    except Exception as e:
        return {"status": "failure", "error": str(e)}